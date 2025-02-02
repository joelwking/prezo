#!/usr/bin/env python
# -*- coding: utf-8 -*-
#
#     Copyright (c) 2019 World Wide Technology, LLC
#     All rights reserved.
#
#     author: joel.king@wwt.com (@joelwking)
#     written:  7 October 2019
#
#     description: Python class to manage and store (Powerpoint) presentations (and other files)
#
import ast
from datetime import timedelta
import datetime
import os
from urllib3.exceptions import ProtocolError
import urllib.parse

from rake_nltk import Rake

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from minio import Minio
from minio.commonconfig import Tags
from minio.error import InvalidResponseError
from minio.error import S3Error
from minio.error import ServerError


class PresentationIndex(object):
    """
        Python class to manage and store (Powerpoint) presentations (and other files)
    """
    KW_NAME = 'rake_keywords'
    DEFAULT = 's3.amazonaws.com'                           # alternately, 'play.min.io'
                                                           # DEFAULT = 'fra1.digitaloceanspaces.com'
    MAX_KEY_LEN = 128                                      # Max tag key length
    MAX_VALUE_LEN = 256                                    # Max tag value length

    def __init__(self, access_key=None, secret_key=None, bucket=None, cloud=DEFAULT):
        """
            metadata is prepended with 'x-amz-meta-' plus the variable name specified in the metadata dictionary

            any error messages are stored in error_message for reference by the calling program
        """
        self.access_key = access_key
        self.secret_key = secret_key
        self.bucket = bucket
        self.minioClient = None
        self.cloud = cloud
        self.error_message = None
        self.message = None
        self.KEYWORDS = 'x-amz-meta-{}'.format(PresentationIndex.KW_NAME)

        self.init_minio()

    def init_minio(self, secure=True):
        """
            Create an instance of the Minio client with the appropriate credentatials
        """

        self.minioClient = Minio(self.cloud, access_key=self.access_key, secret_key=self.secret_key, secure=secure)

    def verify_bucket_exists(self):
        """
            input: self.bucket

            returns: True or False to specify if the bucket exists
        """
        try:
            buckets = self.minioClient.list_buckets()
        except (S3Error, ServerError) as err:
            self.error_message = 'VERIFY_BUCKET_EXISTS:ERROR {}'.format(err)
            return False

        for bucket in buckets:
            if bucket.name == self.bucket:
                self.message = 'VERIFY_BUCKET_EXISTS:INFO ... found bucket {}'.format(bucket.name)
                return True

        return False

    def upload_file(self, filepath=None, metadata=dict(), tags=None, content_type='application/octet-stream'):
        """
            input: filepath: location of the file on the local system
                   metadata: associated meta data
                   content_type: configurable S3 uses the specified default value for .pptx values
 
            returns: None indicating an error, or the etag number of the object

            Filenames may include spaces, thus, urllib.parse.quote_plus
        """

        remote_name = urllib.parse.quote_plus(os.path.basename(filepath))

        try:
            etag = self.minioClient.fput_object(self.bucket, remote_name, filepath, metadata=metadata, tags=tags, content_type=content_type)
        except (InvalidResponseError, S3Error, ServerError, FileNotFoundError) as err:
            self.error_message = err
            return None

        return etag

    def rake_it(self, input_text, depth=10):
        """
            input: depth: maximum number of ranked keyword phrases to return
                   input_text: a list of sentences or text to rake.

            returns: None or a tuple of scores and phrases
        """

        r = Rake()

        if isinstance(input_text, list):
            r.extract_keywords_from_sentences(input_text)            # Extraction given the list of strings where each string is a sentence.
        elif isinstance(input_text, str):
            r.extract_keywords_from_text(input_text)                 # Extraction given the text.
        else:
            self.error_message = 'Input must be either of type list or string'
            return None

        if isinstance(depth, int):
            return r.get_ranked_phrases_with_scores()[0:depth]       # To get keyword phrases ranked highest to lowest with scores.
        else:
            self.error_message = 'depth must be an integer'
            return None

    def extract_text(self, path_to_presentation):
        """
            input: path_to_presentation: filename of the presentation to analyze

            returns: a list of the text found in the presentation

            note: text_runs will be populated with a list of strings, one for each text run in presentation
            reference: https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
        """

        text_runs = []

        prs = self.get_presentation_object(path_to_presentation)

        if not prs:
            return ['error extracting text with pptx']

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return text_runs

    def get_presentation_object(self, path_to_presentation):
        """
            Attempt to read the presentation file and return the object

            input: path_to_presentation: filename of the presentation to analyze

            returns: None if any errors occur, otherwise the presentation object

        """
        try:
            prs = Presentation(path_to_presentation)
        except (KeyError, PackageNotFoundError) as err:
            self.error_message = '{} {}'.format(path_to_presentation, err)
            return None
        except:
            self.error_message = '{} {}'.format(path_to_presentation, 'bare exception')
            return None

        return prs

    def get_core_properties(self, path_to_presentation):
        """"
            Each presentation has values called core_properties which contain meta data like the
            author name, title, revision etc., which in itself, are valuable metadata

            The datatype of these fields are one of integer, string or datetime.

            input: path_to_presentation: filename of the presentation to analyze

            returns: empty dictionary if any errors occur, otherwise a dictionary of properties

            reference: https://python-pptx.readthedocs.io/en/latest/api/presentation.html#coreproperties-objects
        """

        prs = self.get_presentation_object(path_to_presentation)

        if not prs:
            return dict()

        fields = dict(author=prs.core_properties.author,
                      comments=prs.core_properties.comments,
                      category=prs.core_properties.category,
                      subject=prs.core_properties.subject,
                      title=prs.core_properties.title,
                      keywords=prs.core_properties.keywords,
                      revision=prs.core_properties.revision,
                      last_modified_by=prs.core_properties.last_modified_by,
                      created=prs.core_properties.created,
                      identifier=prs.core_properties.identifier,
                      language=prs.core_properties.language,
                      last_printed=prs.core_properties.last_printed,
                      version=prs.core_properties.version,
                      modified=prs.core_properties.modified,
                      content_status=prs.core_properties.content_status 
                    )
        # Remove empty fields, convert datetime to string
        for key, value in list(fields.items()):
            if value in ('', ' ') or value == None:                   
                del fields[key]
            if isinstance(value, datetime.datetime):
                fields[key] = value.strftime("%d-%b-%Y %H:%M:%S")

        return fields

    def get_metadata(self, remote_name):
        """
            Stat the object to return the meta data

            input: remote_name: name of the object within the bucket

            returns: a tuple containing two items, a list of keywords and the stat object
                     containing the fields: 'bucket_name', 'content_type', 'etag', 'is_dir',
                                            'last_modified', 'metadata', 'object_name', 'size'
        """

        keywords = []

        stat = self.minioClient.stat_object(self.bucket, remote_name)
        ## TODO Identify why evaluating the RAKE keywords was necessary
        # try:
        #    keywords = ast.literal_eval(stat.metadata.get(self.KEYWORDS))
        #except (ValueError, SyntaxError) as err:
        #    self.error_message = '{} {}'.format(remote_name, err)
        keywords.append(stat.metadata.get(self.KEYWORDS, ''))
        #
        #  Add non RAKE keywords
        #
        keywords.append(stat.metadata.get('x-amz-meta-author', ''))
        keywords.append(stat.metadata.get('x-amz-meta-subject', ''))
        keywords.append(stat.metadata.get('x-amz-meta-keywords', ''))
        keywords.append(stat.metadata.get('x-amz-meta-title', ''))
        keywords.append(stat.metadata.get('x-amz-meta-last_modified_by', ''))
        keywords.append(stat.object_name)

        return (keywords, stat)

    def us_ascii(self, text):
        """
            Only US-ASCII is permitted as meta-data, we expect a list and return a list
            Remove leading and trailing spaces with strip(), otherwise you may encounter:
            'S3 operation failed; code: SignatureDoesNotMatch'
        """
        ascii_text = []
        for index, value in enumerate(text):
            ascii_text.insert(index, ''.join(i for i in value if ord(i)<128).strip())
        return ascii_text

    def get_download_url(self, remote_name, expires=1):
        """
            Get a URL which can be used to download the file even if the bucket is marked private.

            input: remote_name: name of the object within the bucket
                   self.bucket: name of the bucket

            returns: a URL
        """

        return self.minioClient.presigned_get_object(self.bucket, remote_name, expires=timedelta(days=expires))

    def list_objects(self):
        """
            input: self.bucket: name of the bucket

            returns: a list of objects contained in the bucket
        """

        return self.minioClient.list_objects(self.bucket, recursive=True)

    def set_tags(self, tags, max_tags=10):
        """
            input: dictionary representing tags to associate with the object

            returns: Tags object or None

            You can only have 10 tags per object for S3, we ignore excessive tags or tags with excessive length
        """
        #
        #  TODO, bucket or object, for now, only object
        #
        result = None
        if tags and isinstance(tags, dict):
            result = Tags(for_object=True)
            count = 0
            for key, value in tags.items():
                if len(key) > PresentationIndex.MAX_KEY_LEN or len(value) > PresentationIndex.MAX_VALUE_LEN:
                    self.error_message = 'key or value length exceeds max values'
                    continue
                result[key] = value
                if count >= max_tags:
                    continue
                count += 1
        return result
